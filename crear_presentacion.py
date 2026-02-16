# -*- coding: utf-8 -*-
"""
Generador de Presentación PowerPoint Profesional para TFG
Josu Zabala Muxika - Herramienta Didáctica Interactiva para Estaciones de Bombeo
Estilo: Profesional, fondo blanco, logo UPV/EHU en todas las diapositivas
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Colores del tema profesional
AZUL_EHU = RGBColor(0, 61, 121)           # Azul corporativo UPV/EHU
AZUL_CLARO = RGBColor(59, 142, 208)       # Azul secundario
GRIS_OSCURO = RGBColor(51, 51, 51)        # Texto principal
GRIS_MEDIO = RGBColor(102, 102, 102)      # Texto secundario
GRIS_CLARO = RGBColor(240, 240, 240)      # Fondos suaves
BLANCO = RGBColor(255, 255, 255)
VERDE = RGBColor(46, 125, 50)             # Éxito/Validación
NARANJA = RGBColor(230, 126, 34)          # Destacados
ROJO = RGBColor(192, 57, 43)              # Advertencias

def add_logo(slide, prs, base_dir):
    """Añade el logo de UPV/EHU arriba a la derecha en cada slide"""
    logo_path = os.path.join(base_dir, "assets", "logo_ehu.png")
    if os.path.exists(logo_path):
        # Logo arriba a la derecha
        slide.shapes.add_picture(
            logo_path, 
            prs.slide_width - Inches(2.2), 
            Inches(0.2), 
            width=Inches(2)
        )

def add_footer_line(slide, prs):
    """Añade una línea azul inferior como footer"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, prs.slide_height - Inches(0.15),
        prs.slide_width, Inches(0.15)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_EHU
    line.line.fill.background()

def add_slide_number(slide, prs, num, total):
    """Añade número de diapositiva"""
    num_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        prs.slide_width - Inches(0.8), prs.slide_height - Inches(0.5),
        Inches(0.6), Inches(0.3)
    )
    num_box.fill.background()
    num_box.line.fill.background()
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num}/{total}"
    run.font.size = Pt(10)
    run.font.color.rgb = GRIS_MEDIO

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    base_dir = os.path.dirname(os.path.abspath(__file__))
    total_slides = 12
    
    # ========================================
    # SLIDE 1: PORTADA
    # ========================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo blanco (por defecto)
    
    # Barra superior azul
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = AZUL_EHU
    top_bar.line.fill.background()
    
    # Logo en la barra superior
    logo_path = os.path.join(base_dir, "assets", "logo_ehu.png")
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(0.5), Inches(0.15), width=Inches(2.8))
    
    # Texto en barra superior
    uni_text = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9), Inches(0.3), Inches(4), Inches(0.7))
    uni_text.fill.background()
    uni_text.line.fill.background()
    tf = uni_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "GRADO EN INGENIERÍA MECÁNICA"
    run.font.size = Pt(14)
    run.font.color.rgb = BLANCO
    run.font.bold = True
    
    # Título principal
    title_box = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(12.3), Inches(1.5))
    title_box.fill.background()
    title_box.line.fill.background()
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "DESARROLLO DE UNA HERRAMIENTA DIDÁCTICA INTERACTIVA"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "para la Modelización y Estudio de Estaciones de Bombeo"
    run.font.size = Pt(26)
    run.font.color.rgb = GRIS_OSCURO
    
    # Subtítulo
    sub = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.5))
    sub.fill.background()
    sub.line.fill.background()
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "TRABAJO DE FIN DE GRADO"
    run.font.size = Pt(18)
    run.font.color.rgb = GRIS_MEDIO
    
    # Información del autor
    info_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.8), Inches(6.3), Inches(2))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = GRIS_CLARO
    info_box.line.color.rgb = RGBColor(200, 200, 200)
    
    tf = info_box.text_frame
    tf.word_wrap = True
    
    info_items = [
        ("Alumno:", "Josu Zabala Muxika"),
        ("Tutora:", "Idoya Pellejero Salaverria"),
        ("Departamento:", "Ingeniería Energética"),
        ("Curso:", "2025/2026"),
    ]
    
    for i, (label, value) in enumerate(info_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{label} "
        run.font.size = Pt(14)
        run.font.color.rgb = GRIS_MEDIO
        run = p.add_run()
        run.text = value
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = GRIS_OSCURO
    
    add_footer_line(slide1, prs)
    
    # ========================================
    # SLIDE 2: ÍNDICE
    # ========================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide2, prs, base_dir)
    
    # Título
    title = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(8), Inches(0.8))
    title.fill.background()
    title.line.fill.background()
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Contenido"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    # Línea decorativa
    line = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_CLARO
    line.line.fill.background()
    
    # Items del índice
    items = [
        "1. Introducción y Objetivos",
        "2. Marco Teórico",
        "3. Metodología",
        "4. Problema Nº1: Bombeo entre depósitos",
        "5. Problema Nº2: Fuente de chorro",
        "6. Problema Nº4: Cavitación (NPSH)",
        "7. Validación y Resultados",
        "8. Conclusiones",
    ]
    
    for i, item in enumerate(items):
        y = Inches(1.6 + i * 0.65)
        
        item_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(10), Inches(0.55))
        item_box.fill.background()
        item_box.line.fill.background()
        
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item
        run.font.size = Pt(20)
        run.font.color.rgb = GRIS_OSCURO
    
    add_footer_line(slide2, prs)
    add_slide_number(slide2, prs, 2, total_slides)
    
    # ========================================
    # SLIDE 3: INTRODUCCIÓN Y OBJETIVOS
    # ========================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide3, prs, base_dir)
    
    title = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(8), Inches(0.8))
    title.fill.background()
    title.line.fill.background()
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Introducción y Objetivos"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_CLARO
    line.line.fill.background()
    
    # Motivación
    motiv_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5))
    motiv_box.fill.solid()
    motiv_box.fill.fore_color.rgb = GRIS_CLARO
    motiv_box.line.fill.background()
    
    tf = motiv_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Motivación"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "\nEl aprendizaje de la mecánica de fluidos suele estar ligado a la resolución manual de ejercicios, lo que puede dificultar la comprensión real del comportamiento hidráulico."
    run.font.size = Pt(14)
    run.font.color.rgb = GRIS_OSCURO
    
    # Objetivos
    obj_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5))
    obj_box.fill.solid()
    obj_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    obj_box.line.fill.background()
    
    tf = obj_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Objetivos"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = VERDE
    
    objetivos = [
        "✓ Herramienta didáctica interactiva",
        "✓ Visualización en tiempo real",
        "✓ Validación con error < 2%",
        "✓ Distribución como ejecutable",
    ]
    
    for obj in objetivos:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = obj
        run.font.size = Pt(14)
        run.font.color.rgb = GRIS_OSCURO
    
    # Propuesta
    prop_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5))
    prop_box.fill.solid()
    prop_box.fill.fore_color.rgb = RGBColor(227, 242, 253)
    prop_box.line.fill.background()
    
    tf = prop_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Propuesta"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "\nDesarrollar un software en Python con interfaz gráfica que resuelva y muestre de forma dinámica problemas de estaciones de bombeo del documento IBS (Instalaciones de Bombeo Simples) de Mecánica de Fluidos."
    run.font.size = Pt(16)
    run.font.color.rgb = GRIS_OSCURO
    
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "\n→ El estudiante puede interactuar con los datos y observar cómo varían los resultados en tiempo real."
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = AZUL_CLARO
    
    add_footer_line(slide3, prs)
    add_slide_number(slide3, prs, 3, total_slides)
    
    # ========================================
    # SLIDE 4: MARCO TEÓRICO
    # ========================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide4, prs, base_dir)
    
    title = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(8), Inches(0.8))
    title.fill.background()
    title.line.fill.background()
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Marco Teórico"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    line = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_CLARO
    line.line.fill.background()
    
    # Conceptos en tarjetas
    conceptos = [
        ("Ecuación de Bernoulli", "Conservación de energía por unidad de peso a lo largo de una línea de corriente en régimen estacionario.", RGBColor(227, 242, 253)),
        ("Hazen-Williams", "Fórmula empírica para pérdidas de carga:\nhf = J · L · Q^1.852", RGBColor(232, 245, 233)),
        ("Curvas Características", "CCI (instalación) + CCB (bomba)\n→ Intersección = Punto de funcionamiento", RGBColor(255, 243, 224)),
        ("NPSH y Cavitación", "NPSHdisp ≥ NPSHreq + Margen\nEvita daños por implosión de burbujas", RGBColor(255, 235, 238)),
    ]
    
    for i, (titulo, desc, color) in enumerate(conceptos):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.4)
        y = Inches(1.5 + row * 2.7)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6.1), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        tf = card.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = titulo
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = GRIS_OSCURO
        
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "\n" + desc
        run.font.size = Pt(14)
        run.font.color.rgb = GRIS_OSCURO
    
    add_footer_line(slide4, prs)
    add_slide_number(slide4, prs, 4, total_slides)
    
    # ========================================
    # SLIDE 5: METODOLOGÍA
    # ========================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide5, prs, base_dir)
    
    title = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(8), Inches(0.8))
    title.fill.background()
    title.line.fill.background()
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Metodología y Tecnología"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    line = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_CLARO
    line.line.fill.background()
    
    # Fases
    fases = [
        ("1", "Análisis", "Estudio de problemas del documento IBS"),
        ("2", "Selección", "Elección de ejercicios representativos"),
        ("3", "Implementación", "Motor de cálculo en Python"),
        ("4", "Validación", "Comparación con soluciones teóricas"),
    ]
    
    for i, (num, titulo, desc) in enumerate(fases):
        x = Inches(0.5 + i * 3.2)
        
        # Círculo con número
        circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1), Inches(1.6), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = AZUL_EHU
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = BLANCO
        
        # Texto
        text_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.5), Inches(3), Inches(1.2))
        text_box.fill.background()
        text_box.line.fill.background()
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = titulo
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = GRIS_OSCURO
        
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(12)
        run.font.color.rgb = GRIS_MEDIO
    
    # Tecnologías
    tech_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4), Inches(12.3), Inches(2.8))
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = GRIS_CLARO
    tech_box.line.fill.background()
    
    tf = tech_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Stack Tecnológico"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    techs = [
        ("Python 3.11", "Lenguaje base - NumPy, SciPy"),
        ("CustomTkinter", "Interfaz gráfica moderna"),
        ("Matplotlib", "Visualización de curvas H-Q"),
        ("PyInstaller", "Generación de ejecutable .exe"),
    ]
    
    p = tf.add_paragraph()
    for nombre, desc in techs:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  •  {nombre}: "
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = GRIS_OSCURO
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(14)
        run.font.color.rgb = GRIS_MEDIO
    
    add_footer_line(slide5, prs)
    add_slide_number(slide5, prs, 5, total_slides)
    
    # ========================================
    # SLIDE 6: PROBLEMA 1
    # ========================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide6, prs, base_dir)
    
    title = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    title.fill.background()
    title.line.fill.background()
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Problema Nº1: Bombeo entre Depósitos"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    line = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_CLARO
    line.line.fill.background()
    
    # Descripción
    desc_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(5.5), Inches(2.2))
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(227, 242, 253)
    desc_box.line.fill.background()
    
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Ecuación de la instalación:"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "\nHmi = Δz + k·Q^1.852 + hf,válvula(Q)"
    run.font.size = Pt(16)
    run.font.name = "Consolas"
    run.font.color.rgb = GRIS_OSCURO
    
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "\n• Hazen-Williams para pérdidas\n• Válvula regulable 0-100%\n• Presurización depósito B"
    run.font.size = Pt(13)
    run.font.color.rgb = GRIS_OSCURO
    
    # Modificaciones
    mod_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(5.5), Inches(2.5))
    mod_box.fill.solid()
    mod_box.fill.fore_color.rgb = RGBColor(255, 243, 224)
    mod_box.line.fill.background()
    
    tf = mod_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Mejoras implementadas:"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = NARANJA
    
    mejoras = [
        "✓ Válvula de asiento con Kv del catálogo",
        "✓ Diámetros comerciales (pasos de 25mm)",
        "✓ Modificación de densidad, rugosidad, longitudes",
        "✓ Cálculo de presión límite en depósito B",
    ]
    for m in mejoras:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = m
        run.font.size = Pt(12)
        run.font.color.rgb = GRIS_OSCURO
    
    # Imagen
    img_path = os.path.join(base_dir, "assets", "p1.png")
    if os.path.exists(img_path):
        slide6.shapes.add_picture(img_path, Inches(6.3), Inches(1.3), width=Inches(6.5))
    
    add_footer_line(slide6, prs)
    add_slide_number(slide6, prs, 6, total_slides)
    
    # ========================================
    # SLIDE 7: PROBLEMA 2
    # ========================================
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide7, prs, base_dir)
    
    title = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    title.fill.background()
    title.line.fill.background()
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Problema Nº2: Fuente de Chorro Vertical"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    line = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_CLARO
    line.line.fill.background()
    
    # Imagen a la izquierda
    img_path = os.path.join(base_dir, "assets", "p2.png")
    if os.path.exists(img_path):
        slide7.shapes.add_picture(img_path, Inches(0.3), Inches(1.3), width=Inches(6.5))
    
    # Descripción a la derecha
    desc_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.4), Inches(5.8), Inches(2))
    desc_box.fill.solid()
    desc_box.fill.fore_color.rgb = RGBColor(227, 242, 253)
    desc_box.line.fill.background()
    
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Altura del chorro:"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "\nh = Vc² / (2g)"
    run.font.size = Pt(18)
    run.font.name = "Consolas"
    run.font.color.rgb = GRIS_OSCURO
    
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "\nFamilia de bombas INP 125/250"
    run.font.size = Pt(12)
    run.font.color.rgb = GRIS_MEDIO
    
    # Mejoras
    mod_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(3.6), Inches(5.8), Inches(2.7))
    mod_box.fill.solid()
    mod_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    mod_box.line.fill.background()
    
    tf = mod_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Características del software:"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = VERDE
    
    mejoras = [
        "✓ 5 rodetes disponibles (225-266mm)",
        "✓ Selección automática de bomba óptima",
        "✓ Visualización animada del chorro",
        "✓ Regulación a altura objetivo",
        "✓ Cálculo de costes energéticos",
    ]
    for m in mejoras:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = m
        run.font.size = Pt(12)
        run.font.color.rgb = GRIS_OSCURO
    
    add_footer_line(slide7, prs)
    add_slide_number(slide7, prs, 7, total_slides)
    
    # ========================================
    # SLIDE 8: PROBLEMA 4 (CAVITACIÓN)
    # ========================================
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide8, prs, base_dir)
    
    title = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(10), Inches(0.8))
    title.fill.background()
    title.line.fill.background()
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Problema Nº4: Cavitación (NPSH)"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    line = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_CLARO
    line.line.fill.background()
    
    # Ecuación NPSH
    eq_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1))
    eq_box.fill.solid()
    eq_box.fill.fore_color.rgb = RGBColor(255, 235, 238)
    eq_box.line.fill.background()
    
    tf = eq_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "NPSHdisponible  ≥  NPSHrequerido  +  Margen de seguridad"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = ROJO
    run.font.name = "Consolas"
    
    # Factores
    factores = [
        ("Altitud", "Patm = 10.33 - z/900", "↑ Altitud = ↓ NPSH"),
        ("Temperatura", "Pv(T) interpolada", "↑ Temp = ↓ NPSH"),
        ("Envejecimiento", "hf = k·(1+0.01·años)", "↑ Uso = ↓ NPSH"),
    ]
    
    for i, (titulo, formula, efecto) in enumerate(factores):
        x = Inches(0.5 + i * 4.3)
        
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.6), Inches(4), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = GRIS_CLARO
        card.line.fill.background()
        
        tf = card.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = titulo
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = GRIS_OSCURO
        
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = formula
        run.font.size = Pt(12)
        run.font.name = "Consolas"
        run.font.color.rgb = GRIS_MEDIO
        
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = efecto
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = ROJO
    
    # Imagen
    img_path = os.path.join(base_dir, "assets", "p3.png")
    if os.path.exists(img_path):
        slide8.shapes.add_picture(img_path, Inches(2), Inches(4.4), width=Inches(9.3))
    
    add_footer_line(slide8, prs)
    add_slide_number(slide8, prs, 8, total_slides)
    
    # ========================================
    # SLIDE 9: VALIDACIÓN
    # ========================================
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide9, prs, base_dir)
    
    title = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(8), Inches(0.8))
    title.fill.background()
    title.line.fill.background()
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Validación de Resultados"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    line = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_CLARO
    line.line.fill.background()
    
    # Tabla de resultados
    table_data = [
        ("Magnitud", "Teórico", "Software", "Error"),
        ("Caudal Q", "42.5 L/s", "42.39 L/s", "0.26%"),
        ("Altura H", "31.0 m", "31.04 m", "0.13%"),
        ("Rendimiento η", "76.0%", "76.0%", "0.00%"),
        ("Potencia P", "20.39 kW", "20.35 kW", "0.20%"),
    ]
    
    for row, (c1, c2, c3, c4) in enumerate(table_data):
        y = Inches(1.5 + row * 0.6)
        is_header = row == 0
        
        bg_color = AZUL_EHU if is_header else (GRIS_CLARO if row % 2 == 0 else BLANCO)
        text_color = BLANCO if is_header else GRIS_OSCURO
        
        for col, (text, width) in enumerate([(c1, 2.2), (c2, 2), (c3, 2), (c4, 1.5)]):
            x = Inches(0.5 + sum([2.2, 2, 2, 1.5][:col]))
            cell = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(width), Inches(0.55))
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            cell.line.color.rgb = RGBColor(200, 200, 200)
            
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(14)
            run.font.bold = is_header
            if col == 3 and not is_header:
                run.font.color.rgb = VERDE
                run.font.bold = True
            else:
                run.font.color.rgb = text_color
    
    # Resumen de errores
    summary_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.5), Inches(4.3), Inches(2.5))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    summary_box.line.fill.background()
    
    tf = summary_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Error Máximo"
    run.font.size = Pt(14)
    run.font.color.rgb = GRIS_MEDIO
    
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "< 0.5%"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = VERDE
    
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Problema Nº1"
    run.font.size = Pt(14)
    run.font.color.rgb = GRIS_OSCURO
    
    # Conclusiones de validación
    val_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.2))
    val_box.fill.solid()
    val_box.fill.fore_color.rgb = GRIS_CLARO
    val_box.line.fill.background()
    
    tf = val_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Conclusiones de la validación:"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    conclusiones = [
        "• Errores inferiores al 2% en todos los problemas",
        "• Diferencias atribuibles a interpolación numérica y redondeo",
        "• Comportamiento físico coherente ante valores límite",
        "• Las relaciones de afinidad se aplican correctamente",
    ]
    for c in conclusiones:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = c
        run.font.size = Pt(14)
        run.font.color.rgb = GRIS_OSCURO
    
    add_footer_line(slide9, prs)
    add_slide_number(slide9, prs, 9, total_slides)
    
    # ========================================
    # SLIDE 10: DISCUSIÓN
    # ========================================
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide10, prs, base_dir)
    
    title = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(8), Inches(0.8))
    title.fill.background()
    title.line.fill.background()
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Discusión"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    line = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_CLARO
    line.line.fill.background()
    
    # Fortalezas
    fort_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(2.8))
    fort_box.fill.solid()
    fort_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    fort_box.line.fill.background()
    
    tf = fort_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Fortalezas"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = VERDE
    
    fortalezas = [
        "• Visualización clara del comportamiento hidráulico",
        "• Interactividad que facilita el aprendizaje",
        "• Precisión numérica validada",
        "• Detección automática de condiciones límite",
        "• Distribución como ejecutable autónomo",
    ]
    for f in fortalezas:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f
        run.font.size = Pt(13)
        run.font.color.rgb = GRIS_OSCURO
    
    # Limitaciones
    lim_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6), Inches(2.8))
    lim_box.fill.solid()
    lim_box.fill.fore_color.rgb = RGBColor(255, 243, 224)
    lim_box.line.fill.background()
    
    tf = lim_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Limitaciones"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = NARANJA
    
    limitaciones = [
        "• Solo régimen estacionario",
        "• Problemas limitados al documento IBS",
        "• Curvas de bomba discretas (no experimentales)",
        "• Simplificaciones académicas conscientes",
    ]
    for l in limitaciones:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = l
        run.font.size = Pt(13)
        run.font.color.rgb = GRIS_OSCURO
    
    # Trabajo futuro
    fut_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.2))
    fut_box.fill.solid()
    fut_box.fill.fore_color.rgb = RGBColor(227, 242, 253)
    fut_box.line.fill.background()
    
    tf = fut_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Líneas de trabajo futuro"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    futuro = [
        "• Extensión a régimen transitorio (golpes de ariete)",
        "• Integración con bases de datos de bombas reales",
        "• Análisis energético y optimización de costes",
        "• Aplicación de técnicas de IA para predicción",
    ]
    for f in futuro:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f
        run.font.size = Pt(14)
        run.font.color.rgb = GRIS_OSCURO
    
    add_footer_line(slide10, prs)
    add_slide_number(slide10, prs, 10, total_slides)
    
    # ========================================
    # SLIDE 11: CONCLUSIONES
    # ========================================
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide11, prs, base_dir)
    
    title = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(8), Inches(0.8))
    title.fill.background()
    title.line.fill.background()
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Conclusiones"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    line = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_CLARO
    line.line.fill.background()
    
    conclusiones = [
        ("✓", "Objetivo cumplido", "Se ha digitalizado el documento IBS en una interfaz interactiva que reproduce el razonamiento hidráulico clásico."),
        ("✓", "Precisión validada", "Resultados con errores < 2%, coherentes con las soluciones teóricas."),
        ("✓", "Valor formativo", "Permite experimentar con parámetros y visualizar las consecuencias en tiempo real."),
        ("✓", "Compatibilidad", "La ingeniería clásica es compatible con las tecnologías actuales; la digitalización potencia el criterio técnico."),
    ]
    
    for i, (icon, titulo, desc) in enumerate(conclusiones):
        y = Inches(1.5 + i * 1.4)
        
        card = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.3), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = GRIS_CLARO
        card.line.fill.background()
        
        tf = card.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{icon}  {titulo}"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = VERDE
        
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(14)
        run.font.color.rgb = GRIS_OSCURO
    
    # Cita final
    quote = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6.2), Inches(9.3), Inches(0.9))
    quote.fill.solid()
    quote.fill.fore_color.rgb = RGBColor(227, 242, 253)
    quote.line.fill.background()
    
    tf = quote.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '"La digitalización no sustituye el criterio técnico; lo potencia."'
    run.font.size = Pt(16)
    run.font.italic = True
    run.font.color.rgb = AZUL_EHU
    
    add_footer_line(slide11, prs)
    add_slide_number(slide11, prs, 11, total_slides)
    
    # ========================================
    # SLIDE 12: GRACIAS
    # ========================================
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Barra superior azul
    top_bar = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = AZUL_EHU
    top_bar.line.fill.background()
    
    # Logo
    if os.path.exists(logo_path):
        slide12.shapes.add_picture(logo_path, Inches(0.5), Inches(0.15), width=Inches(2.8))
    
    # Gracias
    thanks = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2))
    thanks.fill.background()
    thanks.line.fill.background()
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "¡Gracias por su atención!"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = AZUL_EHU
    
    # Preguntas
    q = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.7))
    q.fill.background()
    q.line.fill.background()
    tf = q.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "¿Preguntas?"
    run.font.size = Pt(28)
    run.font.color.rgb = GRIS_MEDIO
    
    # Contacto
    contact = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(5), Inches(5.3), Inches(1.5))
    contact.fill.solid()
    contact.fill.fore_color.rgb = GRIS_CLARO
    contact.line.fill.background()
    
    tf = contact.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Josu Zabala Muxika"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = GRIS_OSCURO
    
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "jzabala060@ikasle.ehu.eus"
    run.font.size = Pt(16)
    run.font.color.rgb = AZUL_CLARO
    
    add_footer_line(slide12, prs)
    
    # Guardar
    output_path = os.path.join(base_dir, "TFG_Presentacion_JosuZabala.pptx")
    prs.save(output_path)
    
    print("\n" + "="*60)
    print("✅ PRESENTACIÓN CREADA EXITOSAMENTE")
    print("="*60)
    print(f"\n📁 Archivo: {output_path}")
    print(f"\n📊 Total de diapositivas: {total_slides}")
    print("\n🎯 Contenido:")
    print("    1.  Portada")
    print("    2.  Índice")
    print("    3.  Introducción y Objetivos")
    print("    4.  Marco Teórico")
    print("    5.  Metodología y Tecnología")
    print("    6.  Problema Nº1: Bombeo entre depósitos")
    print("    7.  Problema Nº2: Fuente de chorro")
    print("    8.  Problema Nº4: Cavitación (NPSH)")
    print("    9.  Validación de Resultados")
    print("   10.  Discusión")
    print("   11.  Conclusiones")
    print("   12.  Gracias")
    print("\n💡 Estilo: Profesional, fondo blanco, logo UPV/EHU")
    print("="*60)
    
    return output_path

if __name__ == "__main__":
    create_presentation()
